#!/usr/bin/env python3
"""法律倫理學 AI 申論答題伺服器"""

from flask import Flask, send_from_directory, request, Response, stream_with_context
import anthropic, json, os

BASE_DIR = os.path.dirname(os.path.abspath(__file__))

# ── 本機開發：讀 .env ────────────────────────────────────────
_env_path = os.path.join(BASE_DIR, ".env")
if os.path.exists(_env_path):
    with open(_env_path) as f:
        for line in f:
            line = line.strip()
            if line and not line.startswith("#") and "=" in line:
                k, v = line.split("=", 1)
                if not os.environ.get(k.strip()):
                    os.environ[k.strip()] = v.strip()

app = Flask(__name__, static_folder=BASE_DIR, static_url_path="")

# ── 課程內容 ─────────────────────────────────────────────────
def load_course_material():
    txt_path = os.path.join(BASE_DIR, "course_material.txt")
    if os.path.exists(txt_path):
        with open(txt_path, encoding="utf-8") as f:
            return f.read()
    try:
        from pptx import Presentation
        course_dir = "/Users/syuanwei/Downloads/法律倫理學"
        files = sorted([f for f in os.listdir(course_dir) if f.endswith(".pptx")])
        parts = []
        for fname in files:
            prs = Presentation(os.path.join(course_dir, fname))
            parts.append(f"\n===== {fname} =====")
            for i, slide in enumerate(prs.slides, 1):
                texts = [s.text.strip() for s in slide.shapes
                         if hasattr(s, "text") and s.text.strip()]
                if texts:
                    parts.append(f"[投影片{i}] " + " / ".join(texts))
        return "\n".join(parts)
    except Exception as e:
        return f"（課程資料載入失敗：{e}）"

COURSE_MATERIAL = load_course_material()

def load_laws_reference():
    path = os.path.join(BASE_DIR, "laws_reference.txt")
    if os.path.exists(path):
        with open(path, encoding="utf-8") as f:
            return f.read()
    return ""

LAWS_REFERENCE = load_laws_reference()

# ── 統一 System Prompt ────────────────────────────────────────
SYSTEM_PROMPT = f"""你是台灣律師考試「法律倫理學」科目的專業申論題答題助手，由清華大學蔡采薇律師的課程內容訓練而成。

【絕對限制：唯一知識來源】
你只能引用下方【課程講義】與【法條全文資料庫】中「明確出現過」的法條條號。
以下情況一律禁止，違反即為錯誤答案：
- 禁止引用兩份資料中均未出現過的任何法條（即使你的訓練資料知道該法條存在）
- 禁止用訓練記憶補充任何條文內容，條文全文必須逐字引用【法條全文資料庫】中的原文
- 禁止「延伸引用」資料庫外的相關法條
- 如果兩份資料的法條不足以回答本題，直接說明「本題所需法條在課程資料中未見記載」，不得自行補充

【課程講義】
{COURSE_MATERIAL}

【法條全文資料庫】
{LAWS_REFERENCE}


【寫作要求】
你正在寫台灣律師考試申論題的標準答案，不是教學講義、不是分析報告。

絕對禁止出現以下用語（出現即為不合格）：
「本題核心在於」「本案涉及」「根據課程資料」「值得注意的是」「綜上所述」
「以下分析」「首先說明」「本文認為」「本題考查」「需要特別注意」
任何向讀者解釋「你接下來要做什麼」的句子

正確語氣範例：
✓ 「按律師倫理規範第31條第1項第1款規定：『…』，律師就利害相反之事件不得受任。」
✓ 「查本件甲律師於受任A時，已知悉…，其後復受B委任，顯與前揭規定不符。」
✓ 「從而，甲律師之行為違反律師法第34條，應受申誡處分。」

錯誤語氣範例（禁止）：
✗ 「本題核心在於利益衝突問題的判斷。」
✗ 「根據課程資料，律師倫理規範對此有所規定。」
✗ 「以下分別就大前提、小前提及涵攝加以說明。」

依序輸出（嚴格遵守分隔符，各分隔符單獨佔一行）：

[[[STRUCTURED]]]
條列式三段論，330~480字。
格式：一、二、三、四 分項，各項下可用（一）（二）細分。
直接以引用法條開始第一項，不得有任何前言。

[[[PARAGRAPH]]]
段落式論述（涵攝為核心），330~480字。
一至三段完整段落，不使用條列標題與任何編號。
首句直接引法條，涵攝貫穿全文，末句明確結論。

[[[CONCISE]]]
段落式精簡版，200~299字。
一至兩段，禁止條列符號與編號。
首句法條＋要件，緊接涵攝，末句結論，無多餘文字。

[[[LAWS]]]
本題引用之全部法條，每條格式：

【○○法第X條（第X項）】
條文全文：「…」（逐字引用【法條全文資料庫】原文；資料庫未收錄者寫「請查核全國法規資料庫」）
本題作用：（一句話）

規則：
- 四個區塊引用完全相同的法條
- 法條條號必須在【課程講義】中出現過，課程外法條一律不得引用
- 條文全文只能引用【法條全文資料庫】原文，絕不憑訓練記憶補充
- 使用繁體中文
"""

# ── 路由 ────────────────────────────────────────────────────
@app.route("/")
def index():
    resp = send_from_directory(BASE_DIR, "index.html")
    resp.headers["Cache-Control"] = "no-store"
    return resp

@app.route("/api/essay", methods=["POST"])
def essay():
    data = request.get_json()
    question = (data or {}).get("question", "").strip()
    if not question:
        return {"error": "請輸入題目"}, 400

    client = anthropic.Anthropic(api_key=os.environ.get("ANTHROPIC_API_KEY"))

    def generate():
        try:
            with client.messages.stream(
                model="claude-opus-4-8",
                max_tokens=6000,
                system=SYSTEM_PROMPT,
                messages=[{"role": "user", "content": f"請針對以下申論題作答：\n\n{question}"}]
            ) as stream:
                for text in stream.text_stream:
                    yield f"data: {json.dumps({'text': text}, ensure_ascii=False)}\n\n"
            yield "data: [DONE]\n\n"
        except Exception as e:
            yield f"data: {json.dumps({'error': str(e)}, ensure_ascii=False)}\n\n"

    return Response(
        stream_with_context(generate()),
        mimetype="text/event-stream",
        headers={"Cache-Control": "no-cache", "X-Accel-Buffering": "no"}
    )

if __name__ == "__main__":
    port = int(os.environ.get("PORT", 5001))
    print("✅ 課程資料載入完成")
    print(f"🌐 開啟瀏覽器：http://localhost:{port}")
    app.run(host="0.0.0.0", port=port, debug=False)
